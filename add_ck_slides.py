"""
Mevcut sunum.pptx dosyasına CK Tool açıklama slaytları ekler.
Slayt 3'ten sonra (Yöntem slaytından sonra) 2 yeni slayt eklenir.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import copy
from lxml import etree

DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
VERY_LIGHT = RGBColor(0xF5, 0xF6, 0xFA)
MEDIUM_BLUE = RGBColor(0x00, 0x52, 0x8A)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment
    return txBox


def add_bullets(slide, items, left, top, width, height,
                font_size=16, color=DARK_GRAY, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
    return txBox


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# Mevcut sunumu ac
prs = Presentation("/Users/bilgem/software_design_project/sunum.pptx")

# ================================================================
# SLAYT A: CK Tool Nedir? (Slayt 3'ten sonra = index 3'e ekle)
# ================================================================
slide_layout = prs.slide_layouts[6]  # blank

# Yeni slayt eklemek icin - sonuna ekleyip sonra siralayacagiz
slide_a = prs.slides.add_slide(slide_layout)
add_bg(slide_a, VERY_LIGHT)

# Baslik bari
add_shape_bg(slide_a, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MEDIUM_BLUE)
add_text_box(slide_a, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             "CK Tool Nedir ve Nasil Calistirdik?", font_size=30, color=WHITE, bold=True)

# Sol panel - CK Tool Nedir
add_shape_bg(slide_a, Inches(0.3), Inches(1.3), Inches(6.2), Inches(6), WHITE)
add_text_box(slide_a, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.4),
             "CK Tool (Chidamber-Kemerer Tool)", font_size=20, color=MEDIUM_BLUE, bold=True)

add_bullets(slide_a, [
    "Mauricio Aniche tarafindan gelistirilen acik kaynak Java metrik analiz araci",
    "",
    "Nasil calisir?",
    "  -> Java kaynak kodunu DERLEMEDEN (compile etmeden) analiz eder",
    "  -> .java dosyalarini statik olarak parse eder",
    "  -> Sinif bazinda CSV ciktisi uretir (class.csv)",
    "",
    "Olctugu metrikler:",
    "  -> CBO (Coupling Between Objects)",
    "  -> DIT (Depth of Inheritance Tree)",
    "  -> WMC (Weighted Methods per Class)",
    "  -> RFC (Response for a Class)",
    "  -> LCOM (Lack of Cohesion of Methods)",
    "  -> NOC (Number of Children)",
    "  -> TCC (Tight Class Cohesion)",
    "  -> publicMethodsQty, privateFieldsQty,",
    "     abstractMethodsQty, totalFieldsQty, LOC ...",
], Inches(0.5), Inches(1.9), Inches(5.8), Inches(5.2), font_size=14, color=DARK_GRAY, spacing=Pt(3))

# Sag panel - Nasil kurduk
add_shape_bg(slide_a, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.8), WHITE)
add_text_box(slide_a, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.4),
             "1. Kurulum", font_size=18, color=ACCENT_BLUE, bold=True)

add_text_box(slide_a, Inches(7.0), Inches(1.9), Inches(5.8), Inches(0.4),
             "Maven Central'dan JAR indirildi:", font_size=14, color=DARK_GRAY)

# Kod kutusu - kurulum
add_shape_bg(slide_a, Inches(7.0), Inches(2.4), Inches(5.8), Inches(0.7), RGBColor(0x2D, 0x2D, 0x2D))
add_text_box(slide_a, Inches(7.1), Inches(2.45), Inches(5.6), Inches(0.6),
             "curl -L -o ck-0.7.0-jar-with-dependencies.jar \\\n"
             '  "https://repo1.maven.org/.../ck-0.7.0-..."',
             font_size=11, color=ACCENT_GREEN, bold=False)

# Sag panel alt - Nasil calistirdik
add_shape_bg(slide_a, Inches(6.8), Inches(4.3), Inches(6.2), Inches(3.0), WHITE)
add_text_box(slide_a, Inches(7.0), Inches(4.4), Inches(5.8), Inches(0.4),
             "2. Calistirma Komutu", font_size=18, color=ACCENT_BLUE, bold=True)

add_text_box(slide_a, Inches(7.0), Inches(4.9), Inches(5.8), Inches(0.4),
             "Her surum icin tek komutla:", font_size=14, color=DARK_GRAY)

# Kod kutusu - calistirma
add_shape_bg(slide_a, Inches(7.0), Inches(5.3), Inches(5.8), Inches(1.7), RGBColor(0x2D, 0x2D, 0x2D))
add_text_box(slide_a, Inches(7.1), Inches(5.35), Inches(5.6), Inches(1.6),
             "java -jar ck-0.7.0-jar-with-dependencies.jar \\\n"
             '  "anomaly-detection/src/main/java" \\\n'
             "  false \\     # anonym siniflar\n"
             "  0 \\         # metrik esik\n"
             "  false \\     # JAR dahil mi\n"
             '  "metrics/1.0.0.0/"  # cikti dizini',
             font_size=12, color=ACCENT_GREEN, bold=False)

add_notes(slide_a, """CK TOOL NEDIR (40 sn):
CK Tool, Mauricio Aniche tarafindan gelistirilen acik kaynak bir Java metrik analiz aracidir.

En onemli ozelligi: Java kodunu DERLEMEDEN analiz eder. Yani projeyi build etmenize gerek yok. Dogrudan .java dosyalarini parse edip sinif bazinda CSV ciktisi uretir.

Olctugu metrikler arasinda CBO, DIT, WMC, RFC, LCOM gibi Chidamber-Kemerer metrikleri var. Bunun yaninda TCC, publicMethodsQty, privateFieldsQty gibi ek metrikler de sagliyor ki bunlar QMOOD hesaplamalari icin kritik.

Kurulumu cok basit — Maven Central'dan tek bir JAR dosyasi indiriyorsunuz. Sonra java -jar komutuyla calistiriyorsunuz. Analiz edilecek dizini, cikti dizinini parametre olarak veriyorsunuz.""")


# ================================================================
# SLAYT B: Otomasyon ve CSV Ciktisi
# ================================================================
slide_b = prs.slides.add_slide(slide_layout)
add_bg(slide_b, VERY_LIGHT)

add_shape_bg(slide_b, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MEDIUM_BLUE)
add_text_box(slide_b, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             "CK Tool: Otomasyon ve Ciktilar", font_size=30, color=WHITE, bold=True)

# Sol panel - Otomasyon scripti
add_shape_bg(slide_b, Inches(0.3), Inches(1.3), Inches(6.2), Inches(3.2), WHITE)
add_text_box(slide_b, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.4),
             "3. Otomasyon Scripti (extract_metrics.sh)", font_size=18, color=ACCENT_BLUE, bold=True)

add_text_box(slide_b, Inches(0.5), Inches(1.9), Inches(5.8), Inches(0.3),
             "12 surum icin dongu ile otomatik calistirma:", font_size=14, color=DARK_GRAY)

add_shape_bg(slide_b, Inches(0.5), Inches(2.3), Inches(5.8), Inches(2.0), RGBColor(0x2D, 0x2D, 0x2D))
add_text_box(slide_b, Inches(0.6), Inches(2.35), Inches(5.6), Inches(1.9),
             'VERSIONS=("1.0.0.0" "1.1.0.0" ... "3.0.0.0")\n\n'
             "for VERSION in ${VERSIONS[@]}; do\n"
             '    git checkout "$VERSION"          # surumu sec\n'
             "    java -jar ck-0.7.0-...jar \\     # CK Tool\n"
             '        "src/main/java" ...          # analiz et\n'
             '        "metrics/$VERSION/"           # kaydet\n'
             "done",
             font_size=12, color=ACCENT_GREEN, bold=False)

# Sag panel - Cikti yapisi
add_shape_bg(slide_b, Inches(6.8), Inches(1.3), Inches(6.2), Inches(3.2), WHITE)
add_text_box(slide_b, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.4),
             "4. Cikti Yapisi", font_size=18, color=ACCENT_BLUE, bold=True)

add_shape_bg(slide_b, Inches(7.0), Inches(1.9), Inches(5.8), Inches(2.4), RGBColor(0x2D, 0x2D, 0x2D))
add_text_box(slide_b, Inches(7.1), Inches(1.95), Inches(5.6), Inches(2.3),
             "metrics/\n"
             "|-- 1.0.0.0/class.csv    -> 288 sinif\n"
             "|-- 1.1.0.0/class.csv    -> 327 sinif\n"
             "|-- 1.3.0.0/class.csv    -> 351 sinif\n"
             "|-- 2.0.0.0/class.csv    -> 351 sinif\n"
             "|-- ...                            \n"
             "|-- 2.16.0.0/class.csv   -> 591 sinif (!)\n"
             "|-- 2.19.0.0/class.csv   -> 609 sinif\n"
             "|-- 3.0.0.0/class.csv    -> 609 sinif",
             font_size=13, color=ACCENT_GREEN, bold=False)

# Alt panel - CSV ornegi
add_shape_bg(slide_b, Inches(0.3), Inches(4.7), Inches(12.7), Inches(2.6), WHITE)
add_text_box(slide_b, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.4),
             "5. CSV Ciktisi Nasil Gorunuyor? (class.csv)", font_size=18, color=ACCENT_BLUE, bold=True)

add_text_box(slide_b, Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.3),
             "Her satir bir Java sinifi, her sutun bir metrik:", font_size=14, color=DARK_GRAY)

# CSV header
add_shape_bg(slide_b, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.45), RGBColor(0x2D, 0x2D, 0x2D))
add_text_box(slide_b, Inches(0.6), Inches(5.62), Inches(12.1), Inches(0.4),
             "file, class, type, cbo, wmc, dit, noc, rfc, lcom, tcc, totalMethodsQty, publicMethodsQty, privateFieldsQty, abstractMethodsQty, loc, ...",
             font_size=11, color=ACCENT_ORANGE, bold=False)

# CSV ornek satir
add_shape_bg(slide_b, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.4), RGBColor(0x3D, 0x3D, 0x3D))
add_text_box(slide_b, Inches(0.6), Inches(6.07), Inches(12.1), Inches(0.35),
             "DiscoveryNodeFilterer.java, ...DiscoveryNodeFilterer, class, cbo=15, wmc=42, dit=2, noc=0, rfc=28, lcom=8, tcc=0.3, ..., loc=156",
             font_size=11, color=WHITE, bold=False)

# Aciklama kutulari
metric_boxes = [
    ("cbo=15", "15 sinifa bagli", ACCENT_RED),
    ("wmc=42", "42 karmasiklik", ACCENT_ORANGE),
    ("lcom=8", "Dusuk kohezyon", ACCENT_ORANGE),
    ("privateFieldsQty=6", "6 private alan -> DAM", ACCENT_BLUE),
    ("tcc=0.3", "Sinif ici tutarlilik -> CAM", ACCENT_GREEN),
]
x_pos = 0.5
for metric, desc, color in metric_boxes:
    add_shape_bg(slide_b, Inches(x_pos), Inches(6.55), Inches(2.35), Inches(0.7), color)
    add_text_box(slide_b, Inches(x_pos + 0.05), Inches(6.55), Inches(2.25), Inches(0.3),
                 metric, font_size=12, color=WHITE, bold=True)
    add_text_box(slide_b, Inches(x_pos + 0.05), Inches(6.82), Inches(2.25), Inches(0.3),
                 desc, font_size=11, color=WHITE)
    x_pos += 2.5

add_notes(slide_b, """OTOMASYON VE CIKTILAR (40 sn):
12 surumu tek tek elle calistirmak yerine bir bash scripti yazdik — extract_metrics.sh. Bu script bir dongu icinde her surum icin git checkout yapip CK Tool'u calistiriyor ve ciktilari metrics klasorune kaydediyor.

Cikti olarak her surum icin bir class.csv dosyasi elde ettik. Ornegin 1.0.0.0 surumunde 288 sinif, 2.16.0.0 surumunde 591 sinif analiz edildi.

CSV dosyasinin her satiri bir Java sinifi. Sutunlarda CBO, WMC, DIT, LCOM gibi metrikler var. Ornegin bir sinifin cbo degeri 15 ise, 15 farkli sinifa bagimli demek. wmc 42 ise karmasiklik yuksek.

Bu ham CSV verilerini sonra Python ile okuyup QMOOD tasarim ozelliklerine donusturduk. Ornegin privateFieldsQty bolü totalFieldsQty bize DAM yani Encapsulation degerini veriyor. tcc degeri dogrudan CAM yani Cohesion olarak kullaniliyor.""")


# ================================================================
# Slaytlari yeniden sirala: yeni slaytlar 3. slayttan sonra gelmeli
# Mevcut siralama: 0,1,2,3,4,5,6,7 + 8(yeni_a),9(yeni_b)
# Istenen: 0,1,2,3, 8,9, 4,5,6,7
# ================================================================

from lxml import etree
nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
xml = prs.part._element
sldIdLst = xml.find('.//p:sldIdLst', nsmap)
sldIds = list(sldIdLst)

# Son 2 slayti (index 8,9) al
new_slide_a = sldIds[-2]
new_slide_b = sldIds[-1]

# Kaldir
sldIdLst.remove(new_slide_a)
sldIdLst.remove(new_slide_b)

# 4. pozisyona ekle (index 3'ten sonra)
sldIds_now = list(sldIdLst)
ref_element = sldIds_now[3]  # mevcut 4. slayt (eski index 3)
ref_idx = list(sldIdLst).index(ref_element)
sldIdLst.insert(ref_idx + 1, new_slide_b)
sldIdLst.insert(ref_idx + 1, new_slide_a)

# Kaydet
output_path = "/Users/bilgem/software_design_project/sunum.pptx"
prs.save(output_path)
print(f"Sunum guncellendi: {output_path}")
print(f"Toplam slayt: {len(prs.slides)}")

# Sirayi yazdir
print("\nSlayt sirasi:")
labels = [
    "1. Kapak",
    "2. Proje Ozeti ve Amac",
    "3. Yontem: Analiz Sureci",
    "4. CK Tool Nedir?          <-- YENI",
    "5. CK Tool: Otomasyon      <-- YENI",
    "6. Ham Metrik Sonuclari",
    "7. QMOOD Kalite Nitelikleri",
    "8. Mimari Bozulma + Teknik Borc",
    "9. LLM Karsilastirmasi",
    "10. Sonuc ve Degerlendirme",
]
for label in labels:
    print(f"  {label}")
