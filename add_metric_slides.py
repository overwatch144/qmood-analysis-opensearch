"""
Mevcut sunum.pptx dosyasına CK metrikleri açıklama slaytları ekler.
Slayt 5'ten sonra (CK Tool Otomasyon slaytından sonra) 2 yeni slayt eklenir.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
VERY_LIGHT = RGBColor(0xF5, 0xF6, 0xFA)
MEDIUM_BLUE = RGBColor(0x00, 0x52, 0x8A)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18,
             color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    return txBox


def add_bullets(slide, items, left, top, width, height,
                size=15, color=DARK_GRAY, spacing=Pt(4)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
    return txBox


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# Mevcut sunumu ac
prs = Presentation("/Users/bilgem/software_design_project/sunum.pptx")
blank = prs.slide_layouts[6]


# ================================================================
# SLAYT 1: Temel CK Metrikleri Aciklamasi
# ================================================================
s1 = prs.slides.add_slide(blank)
add_bg(s1, VERY_LIGHT)
add_shape(s1, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MEDIUM_BLUE)
add_text(s1, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
         "CK Metrikleri: Ne Olcuyor, Ne Anlama Geliyor?", size=30, color=WHITE, bold=True)

# 6 metrik kutusu - 3x2 grid
metrics = [
    ("CBO", "Coupling Between Objects",
     "Bir sinifin kac farkli sinifa\nbagli oldugunu olcer.",
     "Yuksekse: Degisiklik yapinca\ndiger siniflar da etkilenir.\nBakim zorlasir.",
     ACCENT_RED, "Orn: cbo=15 -> 15 sinifa bagli"),

    ("DIT", "Depth of Inheritance Tree",
     "Kalitim agacinda sinifin kac\nseviye derinlikte oldugunu olcer.",
     "Yuksekse: Ust siniflardan gelen\ndavranislari anlamak zorlasir.",
     ACCENT_BLUE, "Orn: dit=2 -> 2 seviye kalitim"),

    ("WMC", "Weighted Methods per Class",
     "Siniftaki metotlarin toplam\nkarmasikligini olcer.",
     "Yuksekse: Sinif cok is yapiyor,\ntest etmek ve anlamak zor.",
     ACCENT_ORANGE, "Orn: wmc=42 -> yuksek karmasiklik"),

    ("RFC", "Response for a Class",
     "Bir nesneye mesaj gelince\ncalisabilecek toplam metot sayisi.",
     "Yuksekse: Hata ayiklama\nzorlasir, yan etkiler artar.",
     PURPLE, "Orn: rfc=28 -> 28 metot tepki verebilir"),

    ("LCOM", "Lack of Cohesion of Methods",
     "Sinif icindeki metotlarin\nbirbirleriyle ne kadar ilgisiz oldugu.",
     "Yuksekse: Sinif birden fazla is\nyapiyor, bolunmeli (SRP ihlali).",
     ACCENT_RED, "Orn: lcom=8 -> dusuk kohezyon"),

    ("NOC", "Number of Children",
     "Bir sinifin kac dogrudan\nalt sinifi oldugunu olcer.",
     "Yuksekse: Bu sinif degisirse\ntum cocuklar etkilenir.",
     ACCENT_GREEN, "Orn: noc=3 -> 3 alt sinif var"),
]

for i, (abbr, full, desc, impact, color, example) in enumerate(metrics):
    col = i % 3
    row = i // 3
    left = 0.3 + col * 4.3
    top = 1.3 + row * 3.05

    # Ana kutu
    add_shape(s1, Inches(left), Inches(top), Inches(4.05), Inches(2.85), WHITE)

    # Baslik bari
    add_shape(s1, Inches(left), Inches(top), Inches(4.05), Inches(0.55), color)
    add_text(s1, Inches(left + 0.1), Inches(top + 0.05), Inches(3.85), Inches(0.45),
             f"{abbr} — {full}", size=14, color=WHITE, bold=True)

    # Aciklama
    add_text(s1, Inches(left + 0.1), Inches(top + 0.6), Inches(1.9), Inches(0.3),
             "Ne olcer?", size=11, color=color, bold=True)
    add_text(s1, Inches(left + 0.1), Inches(top + 0.8), Inches(1.9), Inches(0.8),
             desc, size=11, color=DARK_GRAY)

    # Etki
    add_text(s1, Inches(left + 2.05), Inches(top + 0.6), Inches(1.9), Inches(0.3),
             "Yuksekse?", size=11, color=ACCENT_RED, bold=True)
    add_text(s1, Inches(left + 2.05), Inches(top + 0.8), Inches(1.9), Inches(0.8),
             impact, size=11, color=DARK_GRAY)

    # Ornek
    add_shape(s1, Inches(left + 0.05), Inches(top + 2.15), Inches(3.95), Inches(0.55),
              RGBColor(0xF0, 0xF0, 0xF0))
    add_text(s1, Inches(left + 0.15), Inches(top + 2.2), Inches(3.75), Inches(0.45),
             example, size=11, color=color, bold=True)

add_notes(s1, """CK METRIKLERI ACIKLAMASI (50 sn):
CK Tool'un olctugu 6 temel metrigi aciklayalim:

CBO — Coupling Between Objects: Bir sinifin kac farkli sinifa bagli oldugunu olcer. Ornegin cbo=15 ise, bu sinif 15 farkli sinifa bagimli. Yuksek CBO'da bir degisiklik yapinca diger siniflar da etkilenir.

DIT — Depth of Inheritance Tree: Kalitim hiyerarsisindeki derinlik. dit=2 ise 2 seviye kalitim var. Cok derinde anlamak zorlasir.

WMC — Weighted Methods per Class: Siniftaki tum metotlarin toplam karmasikligi. wmc=42 gibi yuksek degerler sinifin cok is yaptigini, test etmenin zor oldugunu gosterir.

RFC — Response for a Class: Bir mesaj geldiginde calisabilecek toplam metot sayisi — sinifin kendi metotlari ve cagirdiklari. Yuksekse yan etkiler artar.

LCOM — Lack of Cohesion of Methods: Metotlarin birbirleriyle ne kadar ilgisiz oldugu. Yuksekse sinif birden fazla sorumluluk tasiyor demektir — Single Responsibility Principle ihlali.

NOC — Number of Children: Dogrudan alt sinif sayisi. Bu sinif degisirse tum cocuklar etkilenir.""")


# ================================================================
# SLAYT 2: Ek Metrikler + QMOOD Eslemesi + Projeden Ornekler
# ================================================================
s2 = prs.slides.add_slide(blank)
add_bg(s2, VERY_LIGHT)
add_shape(s2, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MEDIUM_BLUE)
add_text(s2, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
         "Ek Metrikler, QMOOD Eslemesi ve Projeden Kritik Degerler", size=28, color=WHITE, bold=True)

# Sol panel - Ek metrikler -> QMOOD
add_shape(s2, Inches(0.3), Inches(1.3), Inches(6.2), Inches(3.3), WHITE)
add_text(s2, Inches(0.5), Inches(1.35), Inches(5.8), Inches(0.35),
         "Ek Metrikler -> QMOOD Tasarim Ozelligi Eslemesi", size=16, color=MEDIUM_BLUE, bold=True)

ek_metrics = [
    ("TCC (Tight Class Cohesion)", "Metotlarin ortak alan kullanim orani", "-> CAM (Cohesion)", ACCENT_GREEN),
    ("publicMethodsQty", "Disariya acik metot sayisi", "-> CIS (Messaging)", ACCENT_BLUE),
    ("privateFieldsQty / totalFieldsQty", "Gizli alan orani (kapsulleme)", "-> DAM (Encapsulation)", PURPLE),
    ("totalFieldsQty - staticFieldsQty", "Statik olmayan alan sayisi", "-> MOA (Composition)", ACCENT_ORANGE),
    ("abstractMethodsQty / totalMethodsQty", "Soyut metot orani", "-> ANA (Abstraction)", ACCENT_RED),
    ("abstractMethodsQty + P(DIT>0)", "Polimorfizm gostergesi", "-> NOP (Polymorphism)", ACCENT_RED),
]

y = 1.8
for metric, desc, mapping, color in ek_metrics:
    add_shape(s2, Inches(0.4), Inches(y), Inches(0.12), Inches(0.35), color)
    add_text(s2, Inches(0.6), Inches(y - 0.02), Inches(2.6), Inches(0.2),
             metric, size=11, color=DARK_GRAY, bold=True)
    add_text(s2, Inches(3.2), Inches(y - 0.02), Inches(1.8), Inches(0.2),
             desc, size=10, color=DARK_GRAY)
    add_text(s2, Inches(5.0), Inches(y - 0.02), Inches(1.4), Inches(0.2),
             mapping, size=11, color=color, bold=True)
    y += 0.4

# Sag panel - Somut ornek
add_shape(s2, Inches(6.8), Inches(1.3), Inches(6.2), Inches(3.3), WHITE)
add_text(s2, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.35),
         "Somut Ornek: Bir Sinifin CK Tool Ciktisi", size=16, color=MEDIUM_BLUE, bold=True)

add_shape(s2, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.65), RGBColor(0x2D, 0x2D, 0x2D))
add_text(s2, Inches(7.1), Inches(1.82), Inches(5.6), Inches(0.6),
         "cbo=15, wmc=42, dit=2, lcom=8, rfc=28, tcc=0.3\n"
         "privateFieldsQty=6, totalFieldsQty=10, publicMethodsQty=8",
         size=12, color=ACCENT_GREEN)

yorumlar = [
    ("cbo=15", "15 farkli sinifa bagli (orta-yuksek coupling)", ACCENT_RED),
    ("wmc=42", "Toplam karmasiklik 42 (yuksek — refactoring adayi)", ACCENT_ORANGE),
    ("dit=2", "2 seviye kalitim (makul)", ACCENT_GREEN),
    ("lcom=8", "Metotlar arasi uyum dusuk (SRP ihlali olabilir)", ACCENT_ORANGE),
    ("tcc=0.3", "Metotlarin %30'u ortak veri kullaniyor (dusuk kohezyon)", ACCENT_ORANGE),
    ("private/total = 6/10", "Kapsulleme %60 -> DAM = 0.60", ACCENT_BLUE),
    ("publicMethods=8", "Disariya 8 metot acik -> CIS = 8", ACCENT_BLUE),
]
y = 2.55
for metric, yorum, color in yorumlar:
    add_shape(s2, Inches(7.0), Inches(y), Inches(2.0), Inches(0.3), color)
    add_text(s2, Inches(7.05), Inches(y + 0.02), Inches(1.9), Inches(0.25),
             metric, size=10, color=WHITE, bold=True)
    add_text(s2, Inches(9.1), Inches(y + 0.02), Inches(3.7), Inches(0.25),
             yorum, size=10, color=DARK_GRAY)
    y += 0.33

# Alt panel - Projeden kritik degerler
add_shape(s2, Inches(0.3), Inches(4.8), Inches(12.7), Inches(2.5), WHITE)
add_text(s2, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.35),
         "Projemizden Kritik Degerler: v1.0.0.0 vs v3.0.0.0", size=17, color=MEDIUM_BLUE, bold=True)

# Tablo basligi
headers = ["Metrik", "v1.0.0.0", "v3.0.0.0", "Degisim", "Yorum"]
col_widths = [2.0, 1.5, 1.5, 1.5, 5.8]
col_starts = [0.5]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

add_shape(s2, Inches(0.4), Inches(5.25), Inches(12.4), Inches(0.35), MEDIUM_BLUE)
for j, h in enumerate(headers):
    add_text(s2, Inches(col_starts[j]), Inches(5.27), Inches(col_widths[j]), Inches(0.3),
             h, size=12, color=WHITE, bold=True)

# Tablo verileri
rows = [
    ("Ort. CBO", "8.58", "10.47", "+%22", "Bagimliliklar surekli artti", ACCENT_ORANGE),
    ("Ort. WMC", "11.69", "12.16", "+%4", "Karmasiklik nispeten stabil", ACCENT_GREEN),
    ("Ort. LCOM", "13.89", "18.57", "+%34", "Kohezyon bozuldu", ACCENT_RED),
    ("Max CBO", "118", "217", "+%84", "God Class — 217 sinifa bagli!", ACCENT_RED),
    ("Max WMC", "103", "196", "+%90", "Asiri karmasik sinif", ACCENT_RED),
    ("Max DIT", "7", "7", "Sabit", "Kalitim kontrol altinda", ACCENT_GREEN),
]

for i, (metric, v1, v3, change, comment, color) in enumerate(rows):
    y_row = 5.65 + i * 0.33
    bg_color = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE
    add_shape(s2, Inches(0.4), Inches(y_row), Inches(12.4), Inches(0.33), bg_color)

    vals = [metric, v1, v3, change, comment]
    for j, val in enumerate(vals):
        is_bold = (j == 0 or j == 3)
        c = color if j == 3 else DARK_GRAY
        add_text(s2, Inches(col_starts[j]), Inches(y_row + 0.03),
                 Inches(col_widths[j]), Inches(0.27),
                 val, size=11, color=c, bold=is_bold)

add_notes(s2, """EK METRIKLER VE PROJEDEN DEGERLER (50 sn):
CK Tool sadece temel 6 metrigi degil, ek metrikleri de olcer. Biz bu ek metrikleri QMOOD tasarim ozelliklerine donusturduk.

Ornegin: TCC yani Tight Class Cohesion, metotlarin ortak alan kullanim oranini olcer — biz bunu CAM yani Cohesion olarak kullandik. privateFieldsQty bolu totalFieldsQty bize DAM yani Encapsulation degerini veriyor.

Sag tarafta somut bir ornek var: cbo=15 olan bir sinif 15 farkli sinifa bagli demek — orta-yuksek coupling. wmc=42 ise yuksek karmasiklik, refactoring adayi. tcc=0.3 ise metotlarin sadece yuzde 30'u ortak veri kullaniyor — dusuk kohezyon.

Alt tabloda projemizden kritik degerler: Ortalama CBO 8.58'den 10.47'ye yuzde 22 artmis — bagimliliklar surekli artti. LCOM yuzde 34 artmis — kohezyon bozulmus. En dikkat cekici: Max CBO 118'den 217'ye cikmis! Bir sinif 217 farkli sinifa bagli — bu acik God Class anti-pattern'i. Max DIT ise 7'de sabit kalmis, kalitim kontrol altinda.""")


# ================================================================
# Slaytlari yeniden sirala
# Mevcut: 0..9 + 10(yeni1), 11(yeni2)
# Istenen: 0..4, 10, 11, 5..9  (yeni slaytlar slayt 5'ten sonra)
# ================================================================
nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
xml = prs.part._element
sldIdLst = xml.find('.//p:sldIdLst', nsmap)
sldIds = list(sldIdLst)

new1 = sldIds[-2]
new2 = sldIds[-1]

sldIdLst.remove(new1)
sldIdLst.remove(new2)

# Slayt 5'ten sonra (index 4'ten sonra = index 5 oncesine)
sldIds_now = list(sldIdLst)
ref = sldIds_now[5]  # eski index 5 = Ham Metrik Sonuclari slaytı
ref_idx = list(sldIdLst).index(ref)
sldIdLst.insert(ref_idx, new2)
sldIdLst.insert(ref_idx, new1)

prs.save("/Users/bilgem/software_design_project/sunum.pptx")
print(f"Sunum guncellendi: {len(prs.slides)} slayt")
print()
print("Slayt sirasi:")
labels = [
    "1.  Kapak",
    "2.  Proje Ozeti ve Amac",
    "3.  Yontem: Analiz Sureci",
    "4.  CK Tool Nedir?",
    "5.  CK Tool: Otomasyon ve Ciktilar",
    "6.  CK Metrikleri Aciklamasi       <-- YENI",
    "7.  Ek Metrikler + QMOOD Eslemesi  <-- YENI",
    "8.  Ham Metrik Sonuclari",
    "9.  QMOOD Kalite Nitelikleri",
    "10. Mimari Bozulma + Teknik Borc",
    "11. LLM Karsilastirmasi",
    "12. Sonuc ve Degerlendirme",
]
for l in labels:
    print(f"  {l}")
