"""
Mevcut sunum.pptx'teki Slayt 6 ve Slayt 7'yi guncelleyerek
eksik metrikleri ve tasarim ozelliklerini ekler.
Slayt 6: NOC kaldirilip NOM ve NOA eklenir (8 metrik, 4x2 grid)
Slayt 7: DSC, DCC, MFA, NOM(Complexity) eklenir
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
VERY_LIGHT = RGBColor(0xF5, 0xF6, 0xFA)
MEDIUM_BLUE = RGBColor(0x00, 0x52, 0x8A)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
TEAL = RGBColor(0x00, 0x89, 0x7B)


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18,
             color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    return txBox


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def clear_slide(slide):
    """Slayttaki tum shape'leri siler."""
    sp_list = list(slide.shapes)
    for sp in sp_list:
        sp_elem = sp._element
        sp_elem.getparent().remove(sp_elem)


prs = Presentation("/Users/bilgem/software_design_project/sunum.pptx")

# ================================================================
# SLAYT 6'YI (index 5) YENIDEN OLUSTUR — 8 Ham Metrik (4x2 grid)
# ================================================================
s6 = prs.slides[5]
clear_slide(s6)
add_bg(s6, VERY_LIGHT)

add_shape(s6, Inches(0), Inches(0), Inches(13.333), Inches(1.0), MEDIUM_BLUE)
add_text(s6, Inches(0.8), Inches(0.15), Inches(11), Inches(0.65),
         "Ham Metrikler: CBO, DIT, WMC, RFC, LCOM, NOM, NOA", size=28, color=WHITE, bold=True)

# 7 metrik + NOC (toplam 8) — 4x2 grid
metrics = [
    ("CBO", "Coupling Between Objects",
     "Bir sinifin kac farkli sinifa\nbagli oldugunu olcer.",
     "Yuksekse: Degisiklik yapinca\ndiger siniflar etkilenir.",
     ACCENT_RED, "Projede: 8.58 -> 10.47 (+%22)"),

    ("DIT", "Depth of Inheritance Tree",
     "Kalitim agacinda sinifin kac\nseviye derinlikte oldugu.",
     "Yuksekse: Ust siniflardan gelen\ndavranislari anlamak zorlasir.",
     ACCENT_BLUE, "Projede: Max DIT = 7 (sabit)"),

    ("WMC", "Weighted Methods per Class",
     "Siniftaki metotlarin toplam\nkarmasikligi.",
     "Yuksekse: Sinif cok is yapiyor,\ntest etmek ve anlamak zor.",
     ACCENT_ORANGE, "Projede: 11.69 -> 12.16 (stabil)"),

    ("RFC", "Response for a Class",
     "Mesaj gelince calisabilecek\ntoplam metot sayisi.",
     "Yuksekse: Hata ayiklama\nzorlasir, yan etkiler artar.",
     PURPLE, "Projede: 16.16 -> 15.52"),

    ("LCOM", "Lack of Cohesion of Methods",
     "Metotlarin birbirleriyle ne\nkadar ilgisiz oldugu.",
     "Yuksekse: Sinif birden fazla is\nyapiyor (SRP ihlali).",
     ACCENT_RED, "Projede: 13.89 -> 18.57 (+%34)"),

    ("NOM", "Number of Methods",
     "Siniftaki toplam metot\nsayisini olcer.",
     "Yuksekse: Sinif cok fazla\nsorumluluk tasiyor olabilir.",
     TEAL, "Projede: ort. 6.01 -> 5.69"),

    ("NOA", "Number of Attributes",
     "Siniftaki toplam alan (field)\nsayisini olcer.",
     "Yuksekse: Sinif cok fazla\nveri tasiyor, bolunmeli.",
     TEAL, "Projede: ort. 4.98 -> 4.55"),

    ("NOC", "Number of Children",
     "Bir sinifin kac dogrudan\nalt sinifi oldugu.",
     "Yuksekse: Bu sinif degisirse\ntum cocuklar etkilenir.",
     ACCENT_GREEN, "Projede: ort. 0.04 -> 0.08"),
]

for i, (abbr, full, desc, impact, color, example) in enumerate(metrics):
    col = i % 4
    row = i // 4
    left = 0.2 + col * 3.3
    top = 1.15 + row * 3.1

    # Ana kutu
    add_shape(s6, Inches(left), Inches(top), Inches(3.1), Inches(2.9), WHITE)

    # Baslik bari
    add_shape(s6, Inches(left), Inches(top), Inches(3.1), Inches(0.5), color)
    add_text(s6, Inches(left + 0.08), Inches(top + 0.03), Inches(2.94), Inches(0.44),
             f"{abbr} — {full}", size=12, color=WHITE, bold=True)

    # Ne olcer
    add_text(s6, Inches(left + 0.08), Inches(top + 0.55), Inches(1.45), Inches(0.22),
             "Ne olcer?", size=10, color=color, bold=True)
    add_text(s6, Inches(left + 0.08), Inches(top + 0.75), Inches(1.45), Inches(0.7),
             desc, size=10, color=DARK_GRAY)

    # Yuksekse
    add_text(s6, Inches(left + 1.55), Inches(top + 0.55), Inches(1.45), Inches(0.22),
             "Yuksekse?", size=10, color=ACCENT_RED, bold=True)
    add_text(s6, Inches(left + 1.55), Inches(top + 0.75), Inches(1.45), Inches(0.7),
             impact, size=10, color=DARK_GRAY)

    # Projeden ornek
    add_shape(s6, Inches(left + 0.03), Inches(top + 2.2), Inches(3.04), Inches(0.5),
              RGBColor(0xF0, 0xF0, 0xF0))
    add_text(s6, Inches(left + 0.1), Inches(top + 2.25), Inches(2.9), Inches(0.4),
             example, size=10, color=color, bold=True)

add_notes(s6, """HAM METRIKLER ACIKLAMASI (60 sn):
Projede kullandigimiz 8 ham metrigi aciklayalim:

CBO — Coupling Between Objects: Bir sinifin kac farkli sinifa bagli oldugu. Projemizde ortalama CBO 8.58'den 10.47'ye yukseldi, yuzde 22 artis.

DIT — Depth of Inheritance Tree: Kalitim derinligi. Max DIT tum surumlerde 7'de sabit kaldi, bu olumlu.

WMC — Weighted Methods per Class: Metotlarin toplam karmasikligi. Nispeten stabil kaldi.

RFC — Response for a Class: Mesaj gelince calisabilecek metot sayisi.

LCOM — Lack of Cohesion of Methods: Metotlar arasi uyumsuzluk. 13.89'dan 18.57'ye yukseldi, yuzde 34 artis — siniflar giderek daha fazla sorumluluk tasiyor.

NOM — Number of Methods: Siniftaki toplam metot sayisi. Ortalama 6 civarinda stabil.

NOA — Number of Attributes: Siniftaki alan sayisi. Ortalama 5 civarinda.

NOC — Number of Children: Alt sinif sayisi. Cok dusuk degerler — kalitim az kullanilmis.""")


# ================================================================
# SLAYT 7'YI (index 6) YENIDEN OLUSTUR — 10 Tasarim Ozelligi TAM LISTE
# ================================================================
s7 = prs.slides[6]
clear_slide(s7)
add_bg(s7, VERY_LIGHT)

add_shape(s7, Inches(0), Inches(0), Inches(13.333), Inches(1.0), MEDIUM_BLUE)
add_text(s7, Inches(0.8), Inches(0.15), Inches(11), Inches(0.65),
         "QMOOD: 10 Tasarim Ozelligi ve Hesaplama Yontemi", size=28, color=WHITE, bold=True)

# 10 tasarim ozelligi — 5x2 grid
design_props = [
    ("DSC", "Design Size", "Toplam sinif sayisi",
     "Sistemin buyuklugu", "288 -> 609 (+%111)", ACCENT_BLUE),

    ("ANA", "Abstraction", "abstractMethods / totalMethods\northalamasiyle hesaplanir",
     "Soyutlama duzeyi", "0.0012 -> 0.0112", PURPLE),

    ("DAM", "Encapsulation", "privateFields / totalFields\northalamasiyla hesaplanir",
     "Veri gizleme kalitesi", "0.579 -> 0.449 (-%22!)", ACCENT_RED),

    ("DCC", "Coupling", "Ortalama CBO degeri\nmu(CBO) ile hesaplanir",
     "Siniflar arasi bagimlilik", "8.58 -> 10.47 (+%22)", ACCENT_RED),

    ("CAM", "Cohesion", "TCC (Tight Class Cohesion)\northalamasiyla hesaplanir",
     "Sinif ici tutarlilik", "-0.142 -> -0.138", ACCENT_GREEN),

    ("MOA", "Composition", "totalFields - staticFields\northalamasiyla hesaplanir",
     "Bilesim / alan yogunlugu", "3.08 -> 2.64", ACCENT_ORANGE),

    ("MFA", "Inheritance", "DIT > 0 olan siniflarin\norani ile hesaplanir",
     "Kalitim kullanim orani", "1.00 -> 1.00 (sabit)", ACCENT_BLUE),

    ("NOP", "Polymorphism", "mu(abstractMethods) + P(DIT>0)\nile hesaplanir",
     "Cok bicimlilik duzeyi", "1.003 -> 1.140", PURPLE),

    ("CIS", "Messaging", "Ortalama public metot sayisi\nmu(publicMethods)",
     "Sinif arayuz buyuklugu", "4.48 -> 4.24", TEAL),

    ("NOM", "Complexity", "Ortalama WMC degeri\nmu(WMC) ile hesaplanir",
     "Metot karmasikligi", "11.69 -> 12.16", ACCENT_ORANGE),
]

for i, (abbr, name, calc, meaning, trend, color) in enumerate(design_props):
    col = i % 5
    row = i // 5
    left = 0.15 + col * 2.65
    top = 1.15 + row * 3.1

    # Ana kutu
    add_shape(s7, Inches(left), Inches(top), Inches(2.5), Inches(2.9), WHITE)

    # Baslik
    add_shape(s7, Inches(left), Inches(top), Inches(2.5), Inches(0.55), color)
    add_text(s7, Inches(left + 0.06), Inches(top + 0.02), Inches(2.38), Inches(0.25),
             abbr, size=16, color=WHITE, bold=True)
    add_text(s7, Inches(left + 0.06), Inches(top + 0.27), Inches(2.38), Inches(0.25),
             name, size=11, color=RGBColor(0xE0, 0xE0, 0xE0))

    # Hesaplama
    add_text(s7, Inches(left + 0.06), Inches(top + 0.6), Inches(2.38), Inches(0.2),
             "Hesaplama:", size=9, color=color, bold=True)
    add_text(s7, Inches(left + 0.06), Inches(top + 0.78), Inches(2.38), Inches(0.55),
             calc, size=9, color=DARK_GRAY)

    # Anlami
    add_text(s7, Inches(left + 0.06), Inches(top + 1.35), Inches(2.38), Inches(0.2),
             "Anlami:", size=9, color=color, bold=True)
    add_text(s7, Inches(left + 0.06), Inches(top + 1.53), Inches(2.38), Inches(0.35),
             meaning, size=9, color=DARK_GRAY)

    # Projeden trend
    is_bad = "!" in trend or "+%22" in trend
    trend_color = ACCENT_RED if is_bad else ACCENT_GREEN
    add_shape(s7, Inches(left + 0.03), Inches(top + 2.15), Inches(2.44), Inches(0.55),
              RGBColor(0xF0, 0xF0, 0xF0))
    add_text(s7, Inches(left + 0.06), Inches(top + 2.17), Inches(2.38), Inches(0.2),
             "Projede (v1.0 -> v3.0):", size=8, color=DARK_GRAY, bold=True)
    add_text(s7, Inches(left + 0.06), Inches(top + 2.38), Inches(2.38), Inches(0.25),
             trend, size=10, color=trend_color, bold=True)

add_notes(s7, """10 TASARIM OZELLIGI (60 sn):
QMOOD modelinde ham metriklerden hesaplanan 10 tasarim ozelligi var. Bunlar kalite niteliklerinin girdileridir.

DSC — Design Size: Toplam sinif sayisi. Projemizde 288'den 609'a yuzde 111 buyume.

ANA — Abstraction: Soyut metot orani. Cok dusuk degerler — proje somut sinif agirlikli.

DAM — Encapsulation: Private alan orani, yani kapsulleme. 0.579'dan 0.449'a dustu — yuzde 22 azalma! Bu ciddi bir kapsulleme kaybi.

DCC — Coupling: Ortalama CBO degeri. 8.58'den 10.47'ye yukseldi — bagimliliklar artti.

CAM — Cohesion: TCC ortalamasi. Tum surumlerde negatif — dusuk kohezyon.

MOA — Composition: Statik olmayan alan sayisi ortalamasi.

MFA — Inheritance: DIT buyuk sifir olan siniflarin orani. Sabit kaldi.

NOP — Polymorphism: Soyut metot ortalamasi + kalitim orani. Hafif artis.

CIS — Messaging: Ortalama public metot sayisi. Sinif arayuz buyuklugu.

NOM — Complexity: Ortalama WMC. Metot karmasikligi.

En kritik degisimler: DAM yuzde 22 dustu (kapsulleme kaybi) ve DCC yuzde 22 artti (bagimlilik artisi). Bu ikisi birlikte Flexibility ve Extendibility'nin cokmesine neden oldu.""")


prs.save("/Users/bilgem/software_design_project/sunum.pptx")
print(f"Sunum guncellendi: {len(prs.slides)} slayt")
print()
print("Guncellenen slaytlar:")
print("  Slayt 6: 8 Ham Metrik (CBO, DIT, WMC, RFC, LCOM, NOM, NOA, NOC) — 4x2 grid")
print("  Slayt 7: 10 Tasarim Ozelligi (DSC, ANA, DAM, DCC, CAM, MOA, MFA, NOP, CIS, NOM) — 5x2 grid")
print()
print("Onceki eksikler giderildi:")
print("  + NOM (Number of Methods) eklendi")
print("  + NOA (Number of Attributes) eklendi")
print("  + DSC (Design Size) eklendi")
print("  + DCC (Coupling) eklendi")
print("  + MFA (Inheritance) eklendi")
print("  + NOM/Complexity eklendi")
