"""
Mevcut sunum.pptx'e OpenSearch Anomaly Detection aciklama slaytlari ekler.
Slayt 2'den sonra (Proje Ozeti'nden sonra) 2 yeni slayt eklenir.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
VERY_LIGHT = RGBColor(0xF5, 0xF6, 0xFA)
MEDIUM_BLUE = RGBColor(0x00, 0x52, 0x8A)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
TEAL = RGBColor(0x00, 0x89, 0x7B)
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18,
             color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    return txBox

def add_bullets(slide, items, left, top, width, height,
                size=15, color=DARK_GRAY, spacing=Pt(4)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


prs = Presentation("/Users/bilgem/software_design_project/sunum.pptx")
blank = prs.slide_layouts[6]


# ================================================================
# SLAYT A: OpenSearch Anomaly Detection Nedir?
# ================================================================
sa = prs.slides.add_slide(blank)
add_bg(sa, VERY_LIGHT)

add_shape(sa, Inches(0), Inches(0), Inches(13.333), Inches(1.0), MEDIUM_BLUE)
add_text(sa, Inches(0.8), Inches(0.15), Inches(11), Inches(0.65),
         "OpenSearch Anomaly Detection Nedir?", size=30, color=WHITE, bold=True)

# Sol panel - Tanim
add_shape(sa, Inches(0.3), Inches(1.2), Inches(6.2), Inches(2.8), WHITE)
add_text(sa, Inches(0.5), Inches(1.25), Inches(5.8), Inches(0.35),
         "Proje Tanimi", size=18, color=MEDIUM_BLUE, bold=True)
add_bullets(sa, [
    "OpenSearch: Amazon'un gelistirdigi acik kaynak",
    "arama ve analiz platformu (Elasticsearch forku)",
    "",
    "Anomaly Detection: Bu platform uzerinde calisan",
    "bir eklenti — zaman serisi verilerinde otomatik",
    "olarak anormallikleri (anomalileri) tespit eder.",
    "",
    "Ornek: Saatte 500-700 siparis olan bir sitede",
    "aniden 50'ye dusmesi veya 5000'e firlamasi",
    "-> Anomaly Detection bunu otomatik algilar",
], Inches(0.5), Inches(1.65), Inches(5.8), Inches(2.2),
   size=13, color=DARK_GRAY, spacing=Pt(2))

# Sag panel - Kullanim senaryolari
add_shape(sa, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.8), WHITE)
add_text(sa, Inches(7.0), Inches(1.25), Inches(5.8), Inches(0.35),
         "Gercek Kullanim Senaryolari", size=18, color=MEDIUM_BLUE, bold=True)

scenarios = [
    ("Sunucu Izleme", "CPU, bellek, disk", "CPU aniden %95'e cikmasi", ACCENT_RED),
    ("Ag Guvenligi", "Trafik hacmi", "Gece 3'te anormal artis (DDoS?)", ACCENT_ORANGE),
    ("E-ticaret", "Satis, sayfa goruntulenme", "Odeme sisteminde ani dusus", ACCENT_BLUE),
    ("IoT Sensorleri", "Sicaklik, basinc", "Makinenin anormal titresimi", TEAL),
    ("Uygulama Log", "Hata orani, yanit suresi", "API yanitinin 10x artmasi", PURPLE),
]

y = 1.7
for name, monitor, anomaly, color in scenarios:
    add_shape(sa, Inches(6.9), Inches(y), Inches(0.12), Inches(0.38), color)
    add_text(sa, Inches(7.1), Inches(y - 0.02), Inches(1.6), Inches(0.2),
             name, size=11, color=color, bold=True)
    add_text(sa, Inches(8.7), Inches(y - 0.02), Inches(1.6), Inches(0.2),
             monitor, size=10, color=DARK_GRAY)
    add_text(sa, Inches(10.3), Inches(y - 0.02), Inches(2.5), Inches(0.2),
             anomaly, size=10, color=ACCENT_RED)
    y += 0.42

# Alt panel - Nasil calisir (akis semasi)
add_shape(sa, Inches(0.3), Inches(4.2), Inches(12.7), Inches(3.1), WHITE)
add_text(sa, Inches(0.5), Inches(4.25), Inches(12.3), Inches(0.35),
         "Nasil Calisir? — 4 Adimli Akis", size=18, color=MEDIUM_BLUE, bold=True)

steps = [
    ("1. Veri Toplama", "OpenSearch index'lerinden\nzaman serisi verisi okur\n(orn. her 5 dakikada)", ACCENT_BLUE),
    ("2. Model Egitimi", "Random Cut Forest (RCF)\nalgoritmasi ile verinin\n'normal' davranisini ogrenir", ACCENT_GREEN),
    ("3. Anomali Tespiti", "Yeni veriyi modelle karsilastirir\nAnomali skoru hesaplar (0-10)\nEsik asarsa -> ANOMALI!", ACCENT_ORANGE),
    ("4. Sonuc ve Uyari", "Anomali detaylari kaydedilir\nAlerting eklentisi ile\nbildirim gonderilir", ACCENT_RED),
]

for i, (title, desc, color) in enumerate(steps):
    left = 0.5 + i * 3.15
    add_shape(sa, Inches(left), Inches(4.7), Inches(2.85), Inches(2.3), color)

    add_text(sa, Inches(left + 0.1), Inches(4.75), Inches(2.65), Inches(0.35),
             title, size=14, color=WHITE, bold=True)
    add_text(sa, Inches(left + 0.1), Inches(5.15), Inches(2.65), Inches(1.5),
             desc, size=12, color=WHITE)

    # Oklar (ilk 3 adim arasinda)
    if i < 3:
        add_text(sa, Inches(left + 2.75), Inches(5.4), Inches(0.5), Inches(0.5),
                 "=>", size=24, color=DARK_GRAY, bold=True)

add_notes(sa, """OPENSEARCH ANOMALY DETECTION NEDIR (50 sn):
Projemizin analiz ettigi yazilimi taniyalim.

OpenSearch, Amazon'un gelistirdigi acik kaynak bir arama ve analiz platformudur — Elasticsearch'un forkudur. Anomaly Detection ise bu platformun bir eklentisi — zaman serisi verilerinde otomatik olarak anomalileri tespit eder.

Ornek verelim: Bir e-ticaret sitesinde saatte normalde 500-700 siparis vardir. Aniden 50'ye duserse veya 5000'e firlarsa, Anomaly Detection bunu otomatik algilar ve uyari verir.

Kullanim alanlari cok genis: sunucu izleme, ag guvenligi, IoT sensorleri, uygulama loglari...

4 adimli calisma prensibi var:
1. Veri toplama — OpenSearch indexlerinden periyodik olarak okur
2. Model egitimi — Random Cut Forest algoritmasi ile normalin ne oldugunu ogrenir
3. Anomali tespiti — yeni gelen veriyi modelle karsilastirir, skor hesaplar
4. Sonuc — anomali tespit edilirse uyari gonderir""")


# ================================================================
# SLAYT B: Java Bileşenleri ve Neden QMOOD İçin Uygun
# ================================================================
sb = prs.slides.add_slide(blank)
add_bg(sb, VERY_LIGHT)

add_shape(sb, Inches(0), Inches(0), Inches(13.333), Inches(1.0), MEDIUM_BLUE)
add_text(sb, Inches(0.8), Inches(0.15), Inches(11), Inches(0.65),
         "Java Bilesenleri, RCF Algoritmasi ve Neden Bu Proje?", size=27, color=WHITE, bold=True)

# Sol ust - Java bilesenleri
add_shape(sb, Inches(0.3), Inches(1.2), Inches(6.2), Inches(2.7), WHITE)
add_text(sb, Inches(0.5), Inches(1.25), Inches(5.8), Inches(0.35),
         "Projenin Java Kod Yapisi", size=17, color=MEDIUM_BLUE, bold=True)

components = [
    ("Detector", "Anomali dedektoru tanimlar (hangi index, alan, periyot)", ACCENT_BLUE),
    ("Model (RCF)", "Random Cut Forest makine ogrenmesi modeli", ACCENT_GREEN),
    ("Runner", "Dedektoru periyodik calistirir", ACCENT_ORANGE),
    ("Result", "Tespit edilen anomalileri saklar", PURPLE),
    ("REST API", "Kullanicinin dedektoru olusturma/yonetme API'si", TEAL),
    ("Transport", "OpenSearch dugumleri arasi iletisim", DARK_GRAY),
    ("Feature", "Veriden ozellik cikarma (aggregation)", ACCENT_RED),
]

y = 1.7
for name, desc, color in components:
    add_shape(sb, Inches(0.4), Inches(y), Inches(0.12), Inches(0.28), color)
    add_text(sb, Inches(0.6), Inches(y - 0.02), Inches(1.8), Inches(0.25),
             name, size=12, color=color, bold=True)
    add_text(sb, Inches(2.5), Inches(y - 0.02), Inches(3.9), Inches(0.25),
             desc, size=11, color=DARK_GRAY)
    y += 0.33

# Sag ust - RCF algoritmasi
add_shape(sb, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.7), WHITE)
add_text(sb, Inches(7.0), Inches(1.25), Inches(5.8), Inches(0.35),
         "Random Cut Forest (RCF) Algoritmasi", size=17, color=MEDIUM_BLUE, bold=True)

add_bullets(sb, [
    "Amazon'un gelistirdigi anomali tespit algoritmasi:",
    "",
    "1. Egitim: Gecmis veriden rastgele agaclar olusturur",
    "2. Tahmin: Yeni veriyi agaca eklemeye calisir",
    "3. Skor: Agaci ne kadar degistiriyorsa skor o kadar yuksek",
    "4. Karar: Skor esigi asarsa -> ANOMALI",
], Inches(7.0), Inches(1.65), Inches(5.8), Inches(1.2),
   size=12, color=DARK_GRAY, spacing=Pt(2))

# RCF gorsel aciklama
add_shape(sb, Inches(7.0), Inches(2.85), Inches(5.8), Inches(0.85), RGBColor(0x2D, 0x2D, 0x2D))
add_text(sb, Inches(7.1), Inches(2.87), Inches(5.6), Inches(0.8),
         "Normal:   --*--*--*--*--*--*--     skor: 1-2 (dusuk)\n"
         "Anomali:  --*--*--*--------*--*--  skor: 8-9 (yuksek!)\n"
         "                       ^ ANOMALI NOKTASI",
         size=12, color=ACCENT_GREEN)

# Alt panel - Neden bu proje
add_shape(sb, Inches(0.3), Inches(4.1), Inches(12.7), Inches(1.8), WHITE)
add_text(sb, Inches(0.5), Inches(4.15), Inches(12.3), Inches(0.35),
         "Neden Bu Proje QMOOD Analizi Icin Uygun?", size=17, color=MEDIUM_BLUE, bold=True)

reasons = [
    ("Java Dili", "QMOOD ve CK Tool\nJava icin tasarlanmis", ACCENT_BLUE),
    ("Uygun Boyut", "288-609 sinif\nne kucuk ne buyuk", ACCENT_GREEN),
    ("Zengin OOP", "Kalitim, arayuzler,\nsoyut siniflar yogun", ACCENT_ORANGE),
    ("70+ Surum", "3 major versiyon\n(1.x, 2.x, 3.x)", PURPLE),
    ("Endustriyel", "Amazon/AWS\nuretimde kullaniyor", ACCENT_RED),
]

for i, (title, desc, color) in enumerate(reasons):
    left = 0.4 + i * 2.5
    add_shape(sb, Inches(left), Inches(4.55), Inches(2.3), Inches(1.15), color)
    add_text(sb, Inches(left + 0.05), Inches(4.58), Inches(2.2), Inches(0.3),
             title, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(sb, Inches(left + 0.05), Inches(4.88), Inches(2.2), Inches(0.7),
             desc, size=11, color=WHITE, align=PP_ALIGN.CENTER)

# v2.16 aciklama kutusu
add_shape(sb, Inches(0.3), Inches(6.1), Inches(12.7), Inches(1.15), RGBColor(0xFF, 0xF3, 0xE0))
add_text(sb, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.3),
         "v2.16.0.0'da Ne Oldu? — 227 Yeni Sinifin Nedeni", size=15, color=ACCENT_ORANGE, bold=True)
add_bullets(sb, [
    "-> Yeni anomali tespit modelleri eklendi (forecasting, HC detectors)      "
    "-> Coklu veri akisi destegi (high-cardinality entity detection)      "
    "-> Yeni REST API endpoint'leri",
    "-> Sonuc: Sinif sayisi 364 -> 591 (+%62), ama yeni siniflar daha kucuk ve odakli (WMC dustu, LCOM dustu)",
], Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.7),
   size=12, color=DARK_GRAY, spacing=Pt(3))

add_notes(sb, """JAVA BILESENLERI VE NEDEN BU PROJE (50 sn):
Projenin Java kod yapisi 7 ana bilesenden olusuyor:

Detector — anomali dedektoru tanimlar, hangi veri kaynagi, hangi alan, ne siklikla kontrol edilecek.
Model — Random Cut Forest makine ogrenmesi modeli, projenin kalbi.
Runner — dedektoru periyodik olarak calistirir.
Result — tespit edilen anomalileri saklar.
REST API — kullanicinin dedektoru olusturma ve yonetme arayuzu.
Transport — OpenSearch dugumleri arasi iletisim.
Feature — ham veriden ozellik cikarma.

RCF algoritmasi cok basit bir mantikla calisir: gecmis veriden rastgele agaclar olusturur, yeni gelen veri noktasini agaca eklemeye calisir. Eger bu ekleme agaci cok degistiriyorsa — yani veri beklenmedik bir yerde — anomali skoru yuksek cikar.

Bu proje QMOOD icin ideal cunku: Java dilinde, CK Tool uyumlu. Boyutu makul — 288 ile 609 sinif arasi. OOP yapisi zengin. 70'ten fazla release tag'i var. Ve en onemlisi v2.16'da buyuk bir yapisal degisiklik var — 227 yeni sinif eklenmis. Bu bize analiz icin mukemmel bir kirilma noktasi sagliyor.""")


# ================================================================
# Slaytlari sirala: yeni 2 slayt, Slayt 2'den sonra (index 1'den sonra)
# ================================================================
nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
xml = prs.part._element
sldIdLst = xml.find('.//p:sldIdLst', nsmap)
sldIds = list(sldIdLst)

new_a = sldIds[-2]
new_b = sldIds[-1]
sldIdLst.remove(new_a)
sldIdLst.remove(new_b)

# Index 1'den sonra (Slayt 2 = Proje Ozeti'nden sonra)
sldIds_now = list(sldIdLst)
ref = sldIds_now[1]  # Proje Ozeti slayt
ref_idx = list(sldIdLst).index(ref)
sldIdLst.insert(ref_idx + 1, new_b)
sldIdLst.insert(ref_idx + 1, new_a)


prs.save("/Users/bilgem/software_design_project/sunum.pptx")
print(f"Sunum guncellendi: {len(prs.slides)} slayt")
print()
print("Slayt sirasi:")
labels = [
    "1.  Kapak",
    "2.  Proje Ozeti ve Amac",
    "3.  OpenSearch Anomaly Detection Nedir?     <-- YENI",
    "4.  Java Bilesenleri, RCF, Neden Bu Proje?  <-- YENI",
    "5.  Yontem: Analiz Sureci",
    "6.  CK Tool Nedir?",
    "7.  CK Tool: Otomasyon ve Ciktilar",
    "8.  CK Metrikleri (8 metrik)",
    "9.  QMOOD Tasarim Ozellikleri (10 ozellik)",
    "10. Ham Metrik Sonuclari",
    "11. QMOOD Kalite Nitelikleri Evrimi",
    "12. Mimari Bozulma + Teknik Borc",
    "13. LLM Karsilastirmasi",
    "14. Sonuc ve Degerlendirme",
]
for l in labels:
    print(f"  {l}")
