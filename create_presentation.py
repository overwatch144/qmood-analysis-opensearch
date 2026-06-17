"""
QMOOD Analizi Sunum Oluşturucu
5-6 dakikalık sınıf sunumu için PowerPoint dosyası
Her slaytın altında konuşmacı notları (speaker notes) bulunur.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Renk paleti
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
VERY_LIGHT = RGBColor(0xF5, 0xF6, 0xFA)
MEDIUM_BLUE = RGBColor(0x00, 0x52, 0x8A)

RESULTS_DIR = "/Users/bilgem/software_design_project/results"


def add_bg(slide, color=DARK_BG):
    """Slayta arka plan rengi ekler."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Renkli dikdortgen kutu ekler."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Metin kutusu ekler."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, items, left, top, width, height,
                              font_size=17, color=DARK_GRAY, spacing=Pt(8)):
    """Madde isaretli icerik ekler."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox


def add_notes(slide, text):
    """Konusmaci notlari ekler."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ================================================================
# SLAYT 1: Kapak
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

# Ust dekoratif cizgi
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

# Baslik
add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
             "QMOOD Tabanli Yazilim Kalitesi Analizi\nve LLM Destekli Degerlendirme",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Alt baslik
add_text_box(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8),
             "OpenSearch Anomaly Detection Uzerine Ampirik Bir Calisma",
             font_size=22, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

# Ders bilgisi
add_text_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.6),
             "Yazilim Mimarileri ve Tasarim Yontemleri  |  Donem Projesi",
             font_size=16, color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

# Alt dekoratif cizgi
add_shape_bg(slide, Inches(4), Inches(4.2), Inches(5.3), Inches(0.04), ACCENT_BLUE)

add_notes(slide, """ACILIS (15 sn):
Merhaba, bugunku sunumumda OpenSearch Anomaly Detection projesinin yazilim tasarim kalitesini QMOOD modeli ile analiz ettigimiz ve 3 farkli LLM ile degerlendirdigimiz calismayi sunacagim.

Proje, Java dilinde yazilmis acik kaynak bir anomali tespit eklentisidir. 12 farkli surumunu inceledik.""")


# ================================================================
# SLAYT 2: Proje Ozeti & Amac
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, VERY_LIGHT)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MEDIUM_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             "Proje Ozeti ve Amac", font_size=30, color=WHITE, bold=True)

# Sol kutu - Amac
add_shape_bg(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5), WHITE)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.5),
             "Calismanin Amaci", font_size=20, color=MEDIUM_BLUE, bold=True)
add_bullet_slide_content(slide, [
    "-> Nesne yonelimli yazilim kalitesini olculebilir metriklerle analiz etmek",
    "-> QMOOD modeli ile 12 surumun kalite evrimini izlemek",
    "-> 3 farkli LLM'in kalite degerlendirme yeteneklerini karsilastirmak",
    "-> Prompt muhendisliginin LLM cikti kalitesine etkisini incelemek",
], Inches(0.8), Inches(2.1), Inches(5.2), Inches(4.5), font_size=16)

# Sag kutu - Secilen sistem
add_shape_bg(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.5), WHITE)
add_text_box(slide, Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.5),
             "Secilen Sistem: OpenSearch Anomaly Detection", font_size=18, color=MEDIUM_BLUE, bold=True)
add_bullet_slide_content(slide, [
    "-> Java dilinde, Amazon/AWS destekli",
    "-> 288 - 609 sinif arasi (12 surumde)",
    "-> 16,903 - 35,612 LOC",
    "-> 70+ release tag'i",
    "-> Endustriyel kalitede, uretimde kullanilan",
    "",
    "Secilen 12 Surum (buyuk sicramalarla):",
    "v1.0  ->  v1.1  ->  v1.3  ->  v2.0  ->  v2.2",
    "v2.4  ->  v2.7  ->  v2.10  ->  v2.13",
    "v2.16  ->  v2.19  ->  v3.0",
], Inches(7.1), Inches(2.1), Inches(5.5), Inches(4.5), font_size=15)

add_notes(slide, """PROJE OZETI (45 sn):
Bu calismada OpenSearch Anomaly Detection projesini sectik. Java dilinde yazilmis, Amazon tarafindan desteklenen bir acik kaynak proje.

Neden bu proje? Cunku yeterli buyuklukte — 288'den 609 sinifa buyuyen bir sistem. 70'ten fazla release tag'i var, yani zengin bir evrim sureci mevcut.

12 surum sectik, buyuk sicramalar yaparak — 1.0'dan 3.0'a kadar. Amacimiz QMOOD modeli ile kalite evrimini izlemek ve ardindan ChatGPT, Gemini ve Claude'a bu verileri verip degerlendirmelerini karsilastirmak.""")


# ================================================================
# SLAYT 3: Yontem
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, VERY_LIGHT)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MEDIUM_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             "Yontem: Analiz Sureci", font_size=30, color=WHITE, bold=True)

# Akis seması kutulari
boxes = [
    ("Kaynak Kod\n(12 Surum)", ACCENT_BLUE, 0.3),
    ("CK Tool\n(v0.7.0)", ACCENT_GREEN, 2.85),
    ("Ham Metrikler\n(CBO,DIT,WMC...)", ACCENT_ORANGE, 5.4),
    ("QMOOD\nTasarim Oz.", ACCENT_RED, 7.95),
    ("Kalite\nNitelikleri", RGBColor(0x8E, 0x44, 0xAD), 10.5),
]
for text, color, left in boxes:
    shape = add_shape_bg(slide, Inches(left), Inches(1.5), Inches(2.2), Inches(1.1), color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Calibri'

# Oklar
for x in [2.5, 5.05, 7.6, 10.15]:
    add_text_box(slide, Inches(x), Inches(1.7), Inches(0.4), Inches(0.5),
                 "=>", font_size=22, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# LLM kutusu
shape = add_shape_bg(slide, Inches(4.5), Inches(3.0), Inches(4.3), Inches(0.9), RGBColor(0x2C, 0x3E, 0x50))
tf = shape.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "LLM Degerlendirme: ChatGPT  |  Gemini  |  Claude"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Alt icerik - QMOOD formulleri
add_shape_bg(slide, Inches(0.5), Inches(4.2), Inches(6), Inches(3), WHITE)
add_text_box(slide, Inches(0.7), Inches(4.3), Inches(5.6), Inches(0.4),
             "QMOOD Kalite Nitelikleri (Bansiya & Davis, 2002)", font_size=16, color=MEDIUM_BLUE, bold=True)
add_bullet_slide_content(slide, [
    "Reusability  = -0.25*DCC + 0.25*CAM + 0.5*CIS + 0.5*DSC",
    "Flexibility  = 0.25*DAM - 0.25*DCC + 0.5*MOA + 0.5*NOP",
    "Understandability = -0.33*(ANA+DCC+NOP+NOM+DSC) + 0.33*(DAM+CAM)",
    "Functionality = 0.12*CAM + 0.22*(NOP+CIS+DSC+ANA)",
    "Extendibility = 0.5*(ANA+MFA+NOP) - 0.5*DCC",
    "Effectiveness = 0.2*(ANA+DAM+MOA+MFA+NOP)",
], Inches(0.7), Inches(4.8), Inches(5.6), Inches(2.5), font_size=13, color=DARK_GRAY)

# Sag - Araclar
add_shape_bg(slide, Inches(6.8), Inches(4.2), Inches(6), Inches(3), WHITE)
add_text_box(slide, Inches(7.1), Inches(4.3), Inches(5.5), Inches(0.4),
             "Kullanilan Araclar", font_size=16, color=MEDIUM_BLUE, bold=True)
add_bullet_slide_content(slide, [
    "-> CK Tool v0.7.0 — Java metrik cikarimi (CSV)",
    "-> Python 3.12 + Pandas — QMOOD hesaplama",
    "-> Matplotlib — Gorsellesstirme (6 grafik)",
    "-> Git — Surum yonetimi (git checkout <tag>)",
    "",
    "Metrikler: CBO, DIT, WMC, RFC, LCOM, NOM, NOA",
    "Tasarim Oz: DSC, ANA, DAM, DCC, CAM, MOA, MFA, NOP, CIS, NOM",
], Inches(7.1), Inches(4.8), Inches(5.5), Inches(2.5), font_size=14, color=DARK_GRAY)

add_notes(slide, """YONTEM (50 sn):
Analiz surecimiz su sekilde isledi:
1. Oncelikle repoyu klonladik ve 12 surumu sectik
2. Her surum icin git checkout yapip CK Tool calistirdik — sinif bazli CSV ciktilari elde ettik
3. Bu ham metriklerden (CBO, DIT, WMC, LCOM gibi) QMOOD tasarim ozelliklerini hesapladik
4. Bansiya-Davis formullerini uygulayarak 6 kalite niteligini elde ettik
5. Son olarak tum verileri 3 LLM'e gonderdik

QMOOD modeli uc katmanli: Ham metrikler -> 10 tasarim ozelligi -> 6 kalite niteligini verir. Biz Bansiya ve Davis 2002 formullerini kullandik. Ornegin Reusability, coupling'i negatif, cohesion ve interface boyutunu pozitif agirlikla hesaplaniyor.""")


# ================================================================
# SLAYT 4: Ham Metrik Sonuclari
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, VERY_LIGHT)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MEDIUM_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             "Ham Metrik Sonuclari ve Uc Donem", font_size=30, color=WHITE, bold=True)

# Grafik ekle
img_path = os.path.join(RESULTS_DIR, "raw_metrics_evolution.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.3), Inches(7.5), Inches(5.5))

# Sag panel - 3 donem
add_shape_bg(slide, Inches(8.0), Inches(1.3), Inches(5), Inches(1.5), ACCENT_BLUE)
add_text_box(slide, Inches(8.2), Inches(1.4), Inches(4.6), Inches(0.3),
             "Donem 1: Hizli Buyume (v1.0 - v1.3)", font_size=15, color=WHITE, bold=True)
add_text_box(slide, Inches(8.2), Inches(1.8), Inches(4.6), Inches(0.9),
             "Sinif: 288 -> 351 (+%22)\nLCOM: 13.89 -> 26.50 (2 kat!)\nHizli ozellik ekleme, kohezyon kaybi",
             font_size=13, color=WHITE)

add_shape_bg(slide, Inches(8.0), Inches(3.0), Inches(5), Inches(1.5), ACCENT_GREEN)
add_text_box(slide, Inches(8.2), Inches(3.1), Inches(4.6), Inches(0.3),
             "Donem 2: Stabil Evrim (v2.0 - v2.13)", font_size=15, color=WHITE, bold=True)
add_text_box(slide, Inches(8.2), Inches(3.5), Inches(4.6), Inches(0.9),
             "Sinif: 351 -> 364 (+%3.7)\nMetrikler dar bantta\nOlgun gelistirme donemi",
             font_size=13, color=WHITE)

add_shape_bg(slide, Inches(8.0), Inches(4.7), Inches(5), Inches(1.8), ACCENT_RED)
add_text_box(slide, Inches(8.2), Inches(4.8), Inches(4.6), Inches(0.3),
             "Donem 3: Yapisal Donusum (v2.16 - v3.0)", font_size=15, color=WHITE, bold=True)
add_text_box(slide, Inches(8.2), Inches(5.2), Inches(4.6), Inches(1.2),
             "Sinif: 364 -> 591 (+%62!)\nWMC: 15.30 -> 12.02 (dusus)\nMax CBO: 140 -> 214 (artis)\nBuyuk modul entegrasyonu",
             font_size=13, color=WHITE)

add_notes(slide, """HAM METRIKLER (50 sn):
Soldaki grafiklerde 12 surumun metrik evrimini goruyorsunuz. Uc belirgin donem tespit ettik:

DONEM 1 — Hizli buyume: v1.0'dan v1.3'e sinif sayisi yuzde 22 artti. Ama LCOM — yani kohezyon eksikligi — neredeyse 2 katina cikti. 13.89'dan 26.50'ye. Bu, hizli ozellik eklemenin sinif ici butunlugu bozdugunun net gostergesi.

DONEM 2 — Stabil evrim: v2.0'dan v2.13'e sadece yuzde 3.7 buyume. Metrikler dar bir bantta kalmis. Olgun bir gelistirme donemi.

DONEM 3 — v2.16 kirilma noktasi! Sinif sayisi 364'ten 591'e, yuzde 62 artis. Ilginc olan, ortalama WMC ve LCOM DUSTU. Bu, yeni eklenen 227 sinifin daha basit ve odakli oldugunu gosteriyor. Ancak Max CBO 140'tan 214'e firladi — yani entegrasyon noktalarinda bagimliliklari yogunlasti.""")


# ================================================================
# SLAYT 5: QMOOD Kalite Nitelikleri
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, VERY_LIGHT)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MEDIUM_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             "QMOOD Kalite Nitelikleri Evrimi", font_size=30, color=WHITE, bold=True)

# Kalite evrimi grafigi
img_path = os.path.join(RESULTS_DIR, "quality_attributes_evolution.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.3), Inches(7.8), Inches(5.8))

# Sag panel - ozet bulgular
add_shape_bg(slide, Inches(8.3), Inches(1.3), Inches(4.7), Inches(5.8), WHITE)
add_text_box(slide, Inches(8.5), Inches(1.4), Inches(4.3), Inches(0.4),
             "Temel Bulgular", font_size=18, color=MEDIUM_BLUE, bold=True)

findings = [
    ("Reusability", "4.147 -> 4.093 (%1.3 dusus)", ACCENT_BLUE),
    ("Flexibility", "0.040 -> -0.617 (Pozitiften negatife!)", ACCENT_RED),
    ("Understandability", "Tum surumlerde negatif\nv2.16'da paradoksal iyilesme", ACCENT_ORANGE),
    ("Functionality", "En stabil: 2.988-3.222", ACCENT_GREEN),
    ("Extendibility", "-3.286 -> -4.159 (%26.6 dusus)\nEN KRITIK SORUN", ACCENT_RED),
    ("Effectiveness", "Nispeten stabil: 1.047-1.190", ACCENT_GREEN),
]

y_pos = 1.9
for name, desc, color in findings:
    add_shape_bg(slide, Inches(8.4), Inches(y_pos), Inches(0.15), Inches(0.55), color)
    add_text_box(slide, Inches(8.7), Inches(y_pos), Inches(4.0), Inches(0.25),
                 name, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(8.7), Inches(y_pos + 0.22), Inches(4.0), Inches(0.35),
                 desc, font_size=11, color=DARK_GRAY)
    y_pos += 0.62

add_notes(slide, """KALITE NITELIKLERI (50 sn):
Soldaki grafikte 6 kalite niteliginin 12 surumdeki seyrini goruyorsunuz.

REUSABILITY nispeten stabil — yuzde 1.3 dusus.

FLEXIBILITY en dramatik degisimi gosterdi: v1.0'da pozitif olan tek nitelikti, 0.040. Ama v3.0'da eksi 0.617'ye dustu. Yani sistem esnekligini tamamen kaybetti.

UNDERSTANDABILITY hep negatif — buyuk bir sisteme beklenen bir durum. Ama ilginc bir paradoks var: v2.16'da IYILESTI! Cunku yeni eklenen basit siniflar ortalama karmasikligi dusurdu.

FUNCTIONALITY en stabil nitelik — dar bantta kaldi.

EXTENDIBILITY en kritik sorun — surekli kotulesen trend, yuzde 26.6 dusus. Coupling artisi buna dogrudan neden oluyor.

EFFECTIVENESS nispeten stabil kaldi.""")


# ================================================================
# SLAYT 6: Mimari Bozulma + Radar
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, VERY_LIGHT)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MEDIUM_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             "Mimari Bozulma Analizi ve Teknik Borc", font_size=30, color=WHITE, bold=True)

# Coupling vs Cohesion grafigi
img_path = os.path.join(RESULTS_DIR, "coupling_vs_cohesion.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.3), Inches(5.8), Inches(4.5))

# Radar grafigi
img_path2 = os.path.join(RESULTS_DIR, "quality_radar_comparison.png")
if os.path.exists(img_path2):
    slide.shapes.add_picture(img_path2, Inches(6.3), Inches(1.3), Inches(4.5), Inches(4.5))

# Alt panel - Teknik borc gostergeleri
add_shape_bg(slide, Inches(0.3), Inches(6.0), Inches(12.7), Inches(1.2), RGBColor(0xFF, 0xF3, 0xE0))
add_text_box(slide, Inches(0.5), Inches(6.05), Inches(3), Inches(0.3),
             "Teknik Borc Gostergeleri:", font_size=15, color=ACCENT_ORANGE, bold=True)

indicators = [
    "LCOM ort. 18-27: SRP ihlalleri",
    "Max CBO = 217: God Class",
    "Max WMC = 318: Asiri karmasiklik",
    "Max DIT = 7 (sabit): Kontrol altinda",
]
x_pos = 0.5
for ind in indicators:
    color = ACCENT_RED if "217" in ind or "318" in ind else DARK_GRAY
    add_text_box(slide, Inches(x_pos), Inches(6.4), Inches(3), Inches(0.5),
                 "  " + ind, font_size=12, color=color, bold=("217" in ind or "318" in ind))
    x_pos += 3.2

add_notes(slide, """MIMARI BOZULMA (40 sn):
Sol grafik coupling-cohesion iliskisini gosteriyor. Daire boyutu sinif sayisi. Gordugunuz gibi surumlerde saga kayma var — coupling artiyor ama cohesion degismiyor. Bu klasik mimari bozulma kalıbı.

v2.16 belirgin bir sicrama — buyuk daire, saga kayma.

Sag tarafta radar grafiginde v1.0, v2.7 ve v3.0'in karsilastirmasi var. Flexibility ve Extendibility eksenleri belirgin sekilde kuculmus.

Alt panelde teknik borc gostergeleri: Max CBO 217 — bir sinif 217 farkli sinifa bagli, bu God Class anti-pattern'i. Max WMC 318 — bazi siniflar asiri karmasik. Bunlar acil refactoring gerektiriyor.""")


# ================================================================
# SLAYT 7: LLM Karsilastirmasi
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, VERY_LIGHT)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MEDIUM_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             "LLM Degerlendirme Karsilastirmasi", font_size=30, color=WHITE, bold=True)

# 3 LLM kutusu
llm_data = [
    ("ChatGPT (GPT-4)", ACCENT_BLUE, [
        "Yaklasim: Dengeli, pratik",
        "v2.16: 'Borcun turu degismis'",
        "Teknik borc: Orta",
        "v3.0: 'Kararli durum'",
        "Oneri: DDD, Facade, Event-Driven",
        "Mimari tani: 'Iyi ayristirmis",
        "ama siki entegre'",
    ]),
    ("Google Gemini", ACCENT_ORANGE, [
        "Yaklasim: Elestirel, risk odakli",
        "v2.16: 'Katilasmis (rigid)'",
        "Teknik borc: Yuksek",
        "v3.0: 'Inovasyon kaybi'",
        "Oneri: Hexagonal Architecture",
        "Mimari tani: 'Big Ball of Mud",
        "-> Distributed Spaghetti'",
    ]),
    ("Claude (Anthropic)", RGBColor(0x8E, 0x44, 0xAD), [
        "Yaklasim: Sistematik, nicel",
        "v2.16: 'Paradoksal iyilesme'",
        "Teknik borc: Orta-Yuksek",
        "v3.0: 'Donma noktasi'",
        "Oneri: SRP, Composition > Inheritance",
        "Mimari tani: 'Organik buyume",
        "-> yapisal donusum'",
    ]),
]

for i, (name, color, items) in enumerate(llm_data):
    left = 0.3 + i * 4.35
    add_shape_bg(slide, Inches(left), Inches(1.3), Inches(4.1), Inches(0.6), color)
    add_text_box(slide, Inches(left + 0.1), Inches(1.35), Inches(3.9), Inches(0.5),
                 name, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape_bg(slide, Inches(left), Inches(1.9), Inches(4.1), Inches(3.5), WHITE)
    add_bullet_slide_content(slide, items,
                              Inches(left + 0.15), Inches(2.0), Inches(3.8), Inches(3.3),
                              font_size=13, color=DARK_GRAY, spacing=Pt(4))

# Ortak bulgular
add_shape_bg(slide, Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.7), WHITE)
add_text_box(slide, Inches(0.5), Inches(5.65), Inches(4), Inches(0.3),
             "Uc Modelin Hemfikir Oldugu Tespitler:", font_size=15, color=MEDIUM_BLUE, bold=True)
add_bullet_slide_content(slide, [
    "1. Coupling artisi en kritik risk  |  2. v2.16 kirilma noktasi  |  3. LCOM artisi = SRP ihlali",
    "4. v2.16'da eklenen siniflar daha basit  |  5. God Class ayristirmasi oncelikli refactoring  |  6. v3.0 stabil",
], Inches(0.5), Inches(6.05), Inches(12.3), Inches(1.0), font_size=13, color=DARK_GRAY, spacing=Pt(6))

add_notes(slide, """LLM KARSILASTIRMASI (60 sn):
Uc LLM'e ayni prompt ve metrik verilerini gonderdik.

CHATGPT en dengeli yaklasimi sergiledi. v2.16'yi 'teknik borcun turunun degistigi' bir surum olarak yorumladi — monolitik karmasiklik azalirken entegrasyon borcu artmis. DDD ve Facade pattern onermis.

GEMINI en elestirel tondu. Projenin 'katilastigini' soyledi. En carpici tespiti: mimariyi 'Big Ball of Mud'dan Distributed Spaghetti'ye gecis' olarak tanimladi. v3.0'i 'pazarlama odakli, inovasyon kaybi' seklinde elestirdi. Hexagonal Architecture onerdi.

CLAUDE en sistematik analizi yapti. v2.16'daki LCOM dususunun gercek bir iyilesme degil, yeni eklenen basit siniflarin ortalamaya cekme etkisi oldugunu tespit etti. v3.0'i 'mimari donma noktasi' olarak degerlendirdi.

ORTAK BULGULAR: Uc model de coupling artisini en kritik risk olarak tanimladi, v2.16'yi kirilma noktasi olarak gordu ve God Class ayristirmasini oncelikli buldu.""")


# ================================================================
# SLAYT 8: Sonuc
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.8),
             "Sonuc ve Degerlendirme", font_size=32, color=WHITE, bold=True)

# 5 bulgu kutusu
conclusions = [
    ("1", "Yazilim buyumesi kaliteyi etkiler", "288 -> 609 sinif (%111)\nExtendibility %26.6 dusus", ACCENT_BLUE),
    ("2", "Coupling en kritik sorun", "CBO: 8.58 -> 10.47 (%22.1)\nMax CBO = 217 (God Class)", ACCENT_RED),
    ("3", "v2.16 kirilma noktasi", "227 yeni sinif, WMC dusus\nama DAM %22 azalma", ACCENT_ORANGE),
    ("4", "LLM'ler metrik analizinde kullanisli", "3 model trendleri dogru yorumladi\nAncak baglam eksikligi sinirli", ACCENT_GREEN),
    ("5", "Prompt muhendisligi kritik", "Ham metrik + delta + spesifik soru\n= cikti derinligi artiyor", RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (num, title, desc, color) in enumerate(conclusions):
    left = 0.3 + (i % 5) * 2.6
    top = 1.3

    shape = add_shape_bg(slide, Inches(left), Inches(top), Inches(2.4), Inches(3.0), color)

    # Numara
    add_text_box(slide, Inches(left + 0.1), Inches(top + 0.1), Inches(2.2), Inches(0.5),
                 num, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Baslik
    add_text_box(slide, Inches(left + 0.1), Inches(top + 0.6), Inches(2.2), Inches(0.8),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Aciklama
    add_text_box(slide, Inches(left + 0.1), Inches(top + 1.5), Inches(2.2), Inches(1.2),
                 desc, font_size=11, color=RGBColor(0xE0, 0xE0, 0xE0), alignment=PP_ALIGN.CENTER)

# Gelecek calismalar
add_shape_bg(slide, Inches(0.3), Inches(4.6), Inches(12.7), Inches(2.5), RGBColor(0x25, 0x25, 0x40))
add_text_box(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.4),
             "Gelecek Calismalar", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_slide_content(slide, [
    "-> v3.1-v3.7 surumlerinin dahil edilmesi  |  En yuksek CBO/LCOM siniflarinin tekil analizi",
    "-> Designite veya SonarQube ile capraz dogrulama  |  LLM'lere kaynak kod verilerek spesifik refactoring",
    "-> Baska bir acik kaynak proje ile karsilastirma  |  SQALE modeli ile teknik borc maliyet tahmini",
], Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), font_size=14, color=LIGHT_GRAY, spacing=Pt(8))

# Tesekkur
add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5),
             "Tesekkurler  —  Sorularinizi bekliyorum",
             font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_notes(slide, """SONUC (40 sn):
Ozetlemek gerekirse, 5 temel bulgu:

1. Yazilim buyudukce kalite etkileniyor — proje 2 kattan fazla buyudu, extendibility yuzde 26.6 dustu.

2. Coupling en kritik sorun — CBO monoton artis gosterdi ve Max CBO 217'ye ulasti, bu God Class anti-pattern'i.

3. v2.16 yapisal bir kirilma noktasi — 227 yeni sinif eklendi, WMC dusmus ama kapsulleme bozulmus. ChatGPT'nin dedigi gibi 'teknik borcun turu degismis'.

4. LLM'ler bu tur metrik tabanli analizde gercekten kullanisli — uc model de trendleri dogru tespit etti. Ama projeye ozgu baglam olmadan oneriler genel kaliyor.

5. Prompt muhendisligi cok onemli — ham metrikler, delta verileri ve spesifik sorular ekleyince cikti kalitesi onemli olcude artti.

Gelecekte capraz dogrulama, daha fazla surum analizi ve LLM'lere kaynak kod vererek spesifik refactoring onerileri alinabilir.

Tesekkur ederim, sorularinizi bekliyorum.""")


# ================================================================
# Kaydet
# ================================================================
output_path = "/Users/bilgem/software_design_project/sunum.pptx"
prs.save(output_path)
print(f"Sunum olusturuldu: {output_path}")
print(f"Toplam slayt: {len(prs.slides)}")
